"""Generate presentation.pptx  (8 slides)"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.enum.dml import MSO_THEME_COLOR
import copy
from lxml import etree
import math

# ── Color palette ──────────────────────────────────────
NAVY   = RGBColor(0x1A, 0x2E, 0x48)
BLUE   = RGBColor(0x28, 0x60, 0xA8)
LBLUE  = RGBColor(0x4A, 0x90, 0xC8)
PALE   = RGBColor(0xEE, 0xF5, 0xFC)
ACCENT = RGBColor(0xE8, 0xA0, 0x20)
GRAY   = RGBColor(0x6A, 0x7A, 0x8A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xC8, 0xDC, 0xEA)
SILVER = RGBColor(0xE0, 0xE8, 0xF0)

# ── Slide size: 16:9 widescreen ────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank

# ═══════════════════════════════════════════════════════
# Helper functions
# ═══════════════════════════════════════════════════════

def add_rect(slide, x, y, w, h, fill=None, line=None, radius=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.0)
    else:
        shape.line.fill.background()
    if radius is not None:
        # set rounded corners via XML
        sp = shape._element
        prstGeom = sp.find(qn('p:spPr')).find(qn('a:prstGeom'))
        if prstGeom is not None:
            prstGeom.set('prst', 'roundRect')
            avLst = prstGeom.find(qn('a:avLst'))
            if avLst is None:
                avLst = etree.SubElement(prstGeom, qn('a:avLst'))
            else:
                avLst.clear()
            gd = etree.SubElement(avLst, qn('a:gd'))
            gd.set('name', 'adj')
            gd.set('fmla', f'val {radius}')
    return shape

def add_textbox(slide, x, y, w, h, text, size=12, bold=False,
                color=None, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color if color else NAVY
    return txb

def set_para(para, text, size=12, bold=False, color=None,
             align=PP_ALIGN.LEFT, space_before=0, italic=False):
    para.clear()
    para.alignment = align
    if space_before:
        para.space_before = Pt(space_before)
    run = para.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color if color else NAVY

def add_label(slide, x, y, w, h, label, size=8, bg=BLUE, fg=WHITE):
    shape = add_rect(slide, x, y, w, h, fill=bg, radius=20000)
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = fg
    return shape

def add_header(slide, page_tag, title, right_text=""):
    """Dark navy header bar across top"""
    bar = add_rect(slide, 0, 0, W, Inches(0.85), fill=NAVY)
    # page tag
    add_textbox(slide, Inches(0.4), Inches(0.1), Inches(4), Inches(0.35),
                page_tag, size=8, color=LBLUE)
    # title
    add_textbox(slide, Inches(0.4), Inches(0.4), Inches(9), Inches(0.42),
                title, size=16, bold=False, color=WHITE)
    # right text
    if right_text:
        add_textbox(slide, Inches(10), Inches(0.3), Inches(3.1), Inches(0.4),
                    right_text, size=9, color=RGBColor(0xAA,0xBB,0xCC),
                    align=PP_ALIGN.RIGHT)

def add_card(slide, x, y, w, h, fill=PALE, border=LGRAY):
    return add_rect(slide, x, y, w, h, fill=fill, line=border, radius=40000)

def ph(slide, x, y, w, h=Inches(0.13)):
    """Gray placeholder bar"""
    add_rect(slide, x, y, w, h, fill=SILVER)

def arrow_right(slide, x, y, length, label="", sublabel=""):
    """Horizontal right-pointing arrow line"""
    from pptx.util import Pt
    conn = slide.shapes.add_connector(1, x, y, x+length, y)
    conn.line.color.rgb = BLUE
    conn.line.width = Pt(1.5)
    if label:
        add_textbox(slide, x, y-Inches(0.22), length, Inches(0.22),
                    label, size=8, color=GRAY, align=PP_ALIGN.CENTER)
    if sublabel:
        add_textbox(slide, x, y+Inches(0.04), length, Inches(0.22),
                    sublabel, size=8, color=BLUE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════
# SLIDE 1 – 会社概要
# ═══════════════════════════════════════════════════════
def slide_company_overview():
    slide = prs.slides.add_slide(BLANK)
    add_header(slide, "01 / Company Overview", "会  社  概  要", "SHIKI（親会社）× 日本住戸（グループ会社）")

    top    = Inches(1.0)
    bottom = H - Inches(0.25)
    body_h = bottom - top

    # ── グループ関係 バナー ──────────────────────────────
    banner = add_rect(slide, Inches(0.35), Inches(1.0), W - Inches(0.7), Inches(0.42),
                      fill=NAVY, radius=15000)
    tf = banner.text_frame
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r1 = p.add_run(); r1.text = "SHIKI  は  日本住戸  の"
    r1.font.size = Pt(11); r1.font.color.rgb = RGBColor(0xAA,0xBB,0xCC)
    r2 = p.add_run(); r2.text = "  親会社（出資元）"
    r2.font.size = Pt(11); r2.font.bold = True; r2.font.color.rgb = LBLUE
    r3 = p.add_run(); r3.text = "  であり、同一グループとして協業します"
    r3.font.size = Pt(11); r3.font.color.rgb = RGBColor(0xAA,0xBB,0xCC)

    col_w = (W - Inches(0.35*2 + 0.28)) / 2
    L_x = Inches(0.35)
    R_x = L_x + col_w + Inches(0.28)
    card_top = Inches(1.55)
    card_h   = H - card_top - Inches(0.25)

    # ── SHIKI（親会社）カード ────────────────────────────
    add_card(slide, L_x, card_top, col_w, card_h)
    # 上部アクセントバー
    add_rect(slide, L_x, card_top, col_w, Inches(0.06), fill=BLUE, radius=0)

    tx = L_x + Inches(0.22)
    ty = card_top + Inches(0.18)
    lw = col_w - Inches(0.44)

    # 会社名
    add_textbox(slide, tx, ty, lw, Inches(0.4), "SHIKI", size=20, bold=True, color=BLUE)
    ty += Inches(0.38)
    # 親会社バッジ
    badge = add_label(slide, tx, ty, Inches(1.2), Inches(0.26),
                      "親会社・出資元", size=8, bg=BLUE)
    ty += Inches(0.38)

    rows = [("会社名", "株式会社 SHIKI"), ("所在地", "●●●●●●●●●●"),
            ("設　立", "●●●●年●月"), ("代　表", "●●　●●"),
            ("資本金", "●●●万円"), ("事業内容", "不動産テック・プラットフォーム事業")]
    for lbl, val in rows:
        add_textbox(slide, tx, ty, Inches(1.0), Inches(0.24),
                    lbl, size=8, color=GRAY)
        add_textbox(slide, tx+Inches(1.05), ty, lw-Inches(1.05), Inches(0.24),
                    val, size=10, bold=False, color=NAVY)
        ty += Inches(0.28)

    # ── 日本住戸（グループ会社）カード ──────────────────
    add_card(slide, R_x, card_top, col_w, card_h)
    add_rect(slide, R_x, card_top, col_w, Inches(0.06), fill=NAVY, radius=0)

    tx = R_x + Inches(0.22)
    ty = card_top + Inches(0.18)

    add_textbox(slide, tx, ty, lw, Inches(0.4), "日本住戸", size=20, bold=True, color=NAVY)
    ty += Inches(0.38)
    badge2 = add_label(slide, tx, ty, Inches(1.4), Inches(0.26),
                       "グループ会社（子会社）", size=8, bg=NAVY)
    ty += Inches(0.38)

    rows2 = [("会社名", "●●●●株式会社"), ("所在地", "●●●●●●●●●●"),
             ("設　立", "●●●●年●月"), ("代　表", "●●　●●"),
             ("資本金", "●●●万円"), ("事業内容", "不動産管理・賃貸仲介事業")]
    for lbl, val in rows2:
        add_textbox(slide, tx, ty, Inches(1.0), Inches(0.24),
                    lbl, size=8, color=GRAY)
        add_textbox(slide, tx+Inches(1.05), ty, lw-Inches(1.05), Inches(0.24),
                    val, size=10, color=NAVY)
        ty += Inches(0.28)

    # 出資矢印（カード間）
    mid_x  = L_x + col_w
    arr_cy = card_top + card_h/2
    add_rect(slide, mid_x + Inches(0.04), arr_cy - Inches(0.01),
             Inches(0.20), Pt(1.5), fill=BLUE)
    add_textbox(slide, mid_x, arr_cy - Inches(0.3), Inches(0.28), Inches(0.22),
                "出資", size=7, color=BLUE, align=PP_ALIGN.CENTER)
    add_textbox(slide, mid_x, arr_cy + Inches(0.04), Inches(0.28), Inches(0.22),
                "↓", size=10, color=BLUE, align=PP_ALIGN.CENTER)

    return slide

# ═══════════════════════════════════════════════════════
# SLIDE 2 – ビジネスモデル
# ═══════════════════════════════════════════════════════
def slide_business_model():
    slide = prs.slides.add_slide(BLANK)
    add_header(slide, "02 / Business Model", "ビ ジ ネ ス モ デ ル", "SHIKIモデルを踏襲")

    # 説明バー
    desc = add_card(slide, Inches(0.35), Inches(0.97), W-Inches(0.7), Inches(0.5),
                    fill=PALE, border=LGRAY)
    add_rect(slide, Inches(0.35), Inches(0.97), Inches(0.06), Inches(0.5), fill=BLUE, radius=0)
    add_textbox(slide, Inches(0.55), Inches(1.02), W-Inches(1.0), Inches(0.42),
                "SHIKIが確立した月次課金モデルを基盤に、物件オーナーへのバリュー提供と仲介会社への収益シェアを組み合わせた三方よしのビジネス設計です。",
                size=10, color=NAVY)

    # フロー図
    boxes = [
        ("🏢", "物件オーナー",  "管理コスト削減\n空室リスク低減", False),
        ("⚡", "SHIKI×日本住戸\nPlatform",  "プラットフォーム\n運営・管理",   True),
        ("🏪", "仲介会社",     "手数料収入\nストック収益",  False),
        ("👤", "入 居 者",     "月500円〜\n付帯サービス享受", False),
    ]

    n = len(boxes)
    bw = Inches(2.1)
    bh = Inches(1.7)
    gap = Inches(0.55)
    total_w = n * bw + (n-1) * gap
    sx = (W - total_w) / 2
    by = H - bh - Inches(1.5)

    for i, (icon, title, sub, primary) in enumerate(boxes):
        bx = sx + i*(bw+gap)
        bg = NAVY if primary else PALE
        bc = BLUE if primary else LGRAY
        tc = WHITE if primary else NAVY
        sc = RGBColor(0xAA,0xCC,0xEE) if primary else GRAY

        add_card(slide, bx, by, bw, bh, fill=bg, border=bc)
        add_textbox(slide, bx, by+Inches(0.12), bw, Inches(0.4),
                    icon, size=22, align=PP_ALIGN.CENTER)
        add_textbox(slide, bx, by+Inches(0.5), bw, Inches(0.44),
                    title, size=11, bold=True, color=tc, align=PP_ALIGN.CENTER)
        add_textbox(slide, bx, by+Inches(0.92), bw, Inches(0.66),
                    sub, size=9, color=sc, align=PP_ALIGN.CENTER)

        if i < n-1:
            ax = bx + bw + Inches(0.08)
            ay = by + bh/2
            conn = slide.shapes.add_connector(1, ax, ay, ax+gap-Inches(0.16), ay)
            conn.line.color.rgb = BLUE
            conn.line.width = Pt(1.5)
            lbl = ["委託", "紹介", "マッチング"][i]
            sub_lbl = ["収益還元", "手数料", "月500円〜"][i]
            add_textbox(slide, ax, ay-Inches(0.26), gap-Inches(0.16), Inches(0.22),
                        lbl, size=8, color=GRAY, align=PP_ALIGN.CENTER)
            add_textbox(slide, ax, ay+Inches(0.06), gap-Inches(0.16), Inches(0.22),
                        sub_lbl, size=8, color=BLUE, align=PP_ALIGN.CENTER)

    # 3特徴カード
    features = [
        ("FEATURE 01", "月次サブスクリプション", "入居期間中継続的な収益を確保"),
        ("FEATURE 02", "複数収益源の統合", "電気・コンテンツ・ネットを一本化"),
        ("FEATURE 03", "スケーラブル設計", "戸数増加に応じたストック積み上げ"),
    ]
    fw = (W - Inches(0.7) - Inches(0.28)*2) / 3
    fy = H - Inches(1.2)
    for j, (tag, title, sub) in enumerate(features):
        fx = Inches(0.35) + j*(fw+Inches(0.28))
        add_card(slide, fx, fy, fw, Inches(0.95), fill=PALE, border=LGRAY)
        add_textbox(slide, fx+Inches(0.15), fy+Inches(0.1), fw-Inches(0.3), Inches(0.22),
                    tag, size=7, color=GRAY)
        add_textbox(slide, fx+Inches(0.15), fy+Inches(0.28), fw-Inches(0.3), Inches(0.28),
                    title, size=11, bold=True, color=NAVY)
        add_textbox(slide, fx+Inches(0.15), fy+Inches(0.56), fw-Inches(0.3), Inches(0.3),
                    sub, size=9, color=GRAY)
    return slide

# ═══════════════════════════════════════════════════════
# SLIDE 3 – 仲介ネットワーク
# ═══════════════════════════════════════════════════════
def slide_broker_network():
    slide = prs.slides.add_slide(BLANK)
    add_header(slide, "03 / Broker Network", "仲 介 ネ ッ ト ワ ー ク", "SHIKIネットワークを踏襲")

    body_top = Inches(0.97)
    body_h   = H - body_top - Inches(0.2)
    col_w    = (W - Inches(0.35*2 + 0.28)) / 2
    L_x      = Inches(0.35)
    R_x      = L_x + col_w + Inches(0.28)

    # ── 左: ハブ＋統計 ─────────────────────────────────
    hub_r = Inches(1.05)
    hub_cx = L_x + col_w/2
    hub_cy = body_top + Inches(1.35)
    hub_x  = hub_cx - hub_r
    hub_y  = hub_cy - hub_r

    hub = add_rect(slide, hub_x, hub_y, hub_r*2, hub_r*2, fill=NAVY, line=BLUE, radius=180000)
    add_textbox(slide, hub_x, hub_y+Inches(0.2), hub_r*2, Inches(0.36),
                "🔗", size=20, align=PP_ALIGN.CENTER)
    add_textbox(slide, hub_x, hub_y+Inches(0.55), hub_r*2, Inches(0.28),
                "SHIKI", size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, hub_x, hub_y+Inches(0.8), hub_r*2, Inches(0.22),
                "Platform", size=10, color=LBLUE, align=PP_ALIGN.CENTER)
    add_textbox(slide, hub_x, hub_y+Inches(1.04), hub_r*2, Inches(0.22),
                "全国対応", size=9, color=RGBColor(0xAA,0xCC,0xEE), align=PP_ALIGN.CENTER)

    # 統計グリッド
    stats_y = body_top + Inches(2.85)
    stats = [("●●社", "提携仲介会社数"), ("●●都道府県", "カバーエリア"),
             ("●●万戸", "管理物件数"),  ("●●年", "構築年数")]
    sw = col_w / 2 - Inches(0.1)
    for k, (num, lbl) in enumerate(stats):
        sx = L_x + (k % 2) * (sw + Inches(0.18))
        sy = stats_y + (k // 2) * Inches(0.82)
        add_card(slide, sx, sy, sw, Inches(0.72), fill=PALE, border=LGRAY)
        add_textbox(slide, sx, sy+Inches(0.06), sw, Inches(0.36),
                    num, size=16, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
        add_textbox(slide, sx, sy+Inches(0.42), sw, Inches(0.24),
                    lbl, size=8, color=GRAY, align=PP_ALIGN.CENTER)

    # ── 右: スポーク ────────────────────────────────────
    spokes = [
        ("🏪", "大手仲介チェーン", "全国展開の大手不動産会社と直接提携"),
        ("🏘️", "地域密着型仲介",  "地元に根付いた中小仲介会社ネットワーク"),
        ("💻", "オンライン仲介",  "デジタル仲介プラットフォームとAPI連携"),
        ("🤝", "管理会社直接提携", "物件管理会社を通じた一括導入スキーム"),
    ]
    sh = Inches(0.75)
    s_gap = Inches(0.14)
    total_sh = len(spokes)*sh + (len(spokes)-1)*s_gap
    s_start_y = body_top + (body_h - total_sh) / 2

    for k, (icon, title, sub) in enumerate(spokes):
        sy = s_start_y + k*(sh+s_gap)
        add_card(slide, R_x, sy, col_w, sh, fill=PALE, border=LGRAY)
        add_textbox(slide, R_x+Inches(0.15), sy+Inches(0.15),
                    Inches(0.36), Inches(0.38), icon, size=18)
        add_textbox(slide, R_x+Inches(0.58), sy+Inches(0.1),
                    col_w-Inches(0.72), Inches(0.28),
                    title, size=11, bold=True, color=NAVY)
        add_textbox(slide, R_x+Inches(0.58), sy+Inches(0.38),
                    col_w-Inches(0.72), Inches(0.28),
                    sub, size=9, color=GRAY)

    # 補足メモ
    note_y = s_start_y + total_sh + Inches(0.2)
    note = add_rect(slide, R_x, note_y, col_w, Inches(0.55), fill=NAVY, radius=20000)
    add_textbox(slide, R_x+Inches(0.15), note_y+Inches(0.06),
                col_w-Inches(0.3), Inches(0.44),
                "SHIKIのネットワーク資産をそのまま活用することで、日本住戸は初期コストゼロで全国展開が可能です。",
                size=9, color=RGBColor(0xDD,0xEE,0xFF))

    return slide

# ═══════════════════════════════════════════════════════
# SLIDE 4 – 商品概要
# ═══════════════════════════════════════════════════════
def slide_product():
    slide = prs.slides.add_slide(BLANK)
    add_header(slide, "04 / Product", "商  品  概  要", "料金プラン")

    card_top = Inches(0.97)
    card_h   = H - card_top - Inches(0.2)
    cw = (W - Inches(0.35*2 + 0.28)) / 2
    L_x = Inches(0.35)
    R_x = L_x + cw + Inches(0.28)

    def product_card(x, accent_color, label, price, unit, desc, features):
        add_card(slide, x, card_top, cw, card_h, fill=PALE, border=LGRAY)
        # accent top bar
        add_rect(slide, x, card_top, cw, Inches(0.06), fill=accent_color)
        tx = x + Inches(0.25)
        ty = card_top + Inches(0.2)
        lw = cw - Inches(0.5)
        add_textbox(slide, tx, ty, lw, Inches(0.22), label, size=8, color=GRAY)
        ty += Inches(0.28)
        add_textbox(slide, tx, ty, lw, Inches(0.7), price, size=44, bold=True, color=accent_color)
        ty += Inches(0.62)
        add_textbox(slide, tx, ty, lw, Inches(0.24), unit, size=10, color=GRAY)
        ty += Inches(0.32)
        # divider
        add_rect(slide, tx, ty, lw, Pt(1), fill=LGRAY)
        ty += Inches(0.18)
        add_textbox(slide, tx, ty, lw, Inches(0.55), desc, size=10, color=NAVY)
        ty += Inches(0.62)
        for feat in features:
            add_textbox(slide, tx, ty, Inches(0.2), Inches(0.24), "✓", size=10, bold=True, color=accent_color)
            add_textbox(slide, tx+Inches(0.22), ty, lw-Inches(0.22), Inches(0.24),
                        feat, size=10, color=NAVY)
            ty += Inches(0.28)

    product_card(
        L_x, BLUE,
        "月額基本プラン",
        "¥500",
        "/ 月 · 入居者1名あたり",
        "入居者が毎月負担する基本利用料。\n付帯サービスへのアクセス権が含まれます。",
        ["電気・コンテンツ・ネットが利用可能", "入居期間中継続課金", "初期費用なし"]
    )
    product_card(
        R_x, ACCENT,
        "仲介手数料シェア",
        "  ●%",
        "収益に対する仲介会社への還元率",
        "月額収益の一定割合を仲介会社へ自動還元。\n入居者在籍中はストック収益が積み上がります。",
        ["自動精算・振込", "リアルタイムダッシュボード管理", "戸数に応じた段階レート"]
    )
    return slide

# ═══════════════════════════════════════════════════════
# SLIDE 5 – 付帯サービス
# ═══════════════════════════════════════════════════════
def slide_services():
    slide = prs.slides.add_slide(BLANK)
    add_header(slide, "05 / Add-on Services", "付 帯 サ ー ビ ス", "3つのバリュー")

    card_top = Inches(0.97)
    card_h   = H - card_top - Inches(0.2)
    n = 3
    gap = Inches(0.22)
    cw  = (W - Inches(0.35*2) - gap*(n-1)) / n

    services = [
        ("⚡", "電　　　気", "ENERGY",
         RGBColor(0xFF,0xF3,0xE0), RGBColor(0xE8,0xA0,0x20),
         "入居者向けの電力供給サービス。割安な電気料金プランと使用量の可視化・節電サポートを提供。物件全体での一括契約によりオーナーへのコスト削減効果も実現。",
         "入居者メリット", "●●●●●●●●"),
        ("🎬", "コンテンツ", "CONTENTS",
         RGBColor(0xF3,0xE5,0xF5), RGBColor(0x90,0x60,0xB8),
         "動画・音楽・電子書籍などのデジタルコンテンツを月額サービスに含めて提供。入居の魅力向上と差別化を実現し空室リスクの低減に貢献します。",
         "提供コンテンツ", "●●●●●●●●"),
        ("📡", "インターネット", "NETWORK",
         RGBColor(0xE3,0xF2,0xFD), RGBColor(0x20,0x80,0xC0),
         "物件全体への高速インターネット回線を一括導入。入居者は追加契約不要で即日利用可能。回線品質の管理・保守もプラットフォームが対応します。",
         "回線スペック", "●●●●●●●●"),
    ]

    for i, (icon, name, tag, icon_bg, icon_color, desc, blbl, bval) in enumerate(services):
        cx = Inches(0.35) + i*(cw+gap)
        add_card(slide, cx, card_top, cw, card_h, fill=PALE, border=LGRAY)

        tx = cx + Inches(0.22)
        ty = card_top + Inches(0.2)
        lw = cw - Inches(0.44)

        # icon box
        icon_box = add_rect(slide, tx, ty, Inches(0.65), Inches(0.65),
                            fill=icon_bg, radius=30000)
        add_textbox(slide, tx, ty+Inches(0.1), Inches(0.65), Inches(0.44),
                    icon, size=24, align=PP_ALIGN.CENTER)
        ty += Inches(0.82)

        add_textbox(slide, tx, ty, lw, Inches(0.34),
                    name, size=14, bold=True, color=NAVY)
        ty += Inches(0.38)
        add_label(slide, tx, ty, Inches(1.1), Inches(0.24), tag, size=8)
        ty += Inches(0.38)
        add_textbox(slide, tx, ty, lw, Inches(1.4),
                    desc, size=9, color=GRAY)
        ty += Inches(1.5)

        # benefit box
        bcard = add_rect(slide, tx, ty, lw, Inches(0.65), fill=WHITE, line=LGRAY, radius=15000)
        add_textbox(slide, tx+Inches(0.12), ty+Inches(0.06), lw-Inches(0.24), Inches(0.2),
                    blbl, size=7, color=GRAY)
        add_textbox(slide, tx+Inches(0.12), ty+Inches(0.28), lw-Inches(0.24), Inches(0.28),
                    bval, size=11, bold=True, color=BLUE)

    return slide

# ═══════════════════════════════════════════════════════
# SLIDE 6 – 提案（スキーム図）
# ═══════════════════════════════════════════════════════
def slide_proposal():
    slide = prs.slides.add_slide(BLANK)
    add_header(slide, "06 / Proposal", "提　　　　案", "スキーム図")

    # Note bar at bottom
    note_h  = Inches(0.56)
    note_y  = H - note_h - Inches(0.18)
    note_bg = add_rect(slide, Inches(0.35), note_y, W-Inches(0.7), note_h, fill=NAVY, radius=15000)
    add_textbox(slide, Inches(0.6), note_y+Inches(0.08), W-Inches(1.2), Inches(0.42),
                "提案要旨：SHIKIが構築したプラットフォームと仲介ネットワークを活用し、日本住戸の物件・顧客基盤を組み合わせることで、"
                "入居者への付加価値提供・仲介会社への継続収益・オーナーへのコスト削減を同時に実現します。",
                size=9, color=RGBColor(0xDD,0xEE,0xFF))

    # Scheme diagram area
    diag_top = Inches(0.97)
    diag_h   = note_y - diag_top - Inches(0.14)
    diag_bg  = add_card(slide, Inches(0.35), diag_top, W-Inches(0.7), diag_h, fill=PALE, border=LGRAY)

    # ── Scheme boxes ──────────────────────────────────
    bw = Inches(1.9)
    bh = Inches(1.6)
    positions = [
        (Inches(0.65), (diag_top + diag_h/2) - bh/2),     # オーナー
        (W/2 - bw/2,   (diag_top + diag_h/2) - bh/2),     # Platform
        (W - Inches(0.65) - bw*2 - Inches(0.5),            # 仲介
         (diag_top + diag_h/2) - bh/2),
        (W - Inches(0.65) - bw, (diag_top + diag_h/2) - bh/2),   # 入居者
    ]
    # recompute evenly
    n = 4
    total_bw = n*bw + (n-1)*Inches(0.6)
    sx = (W - total_bw)/2
    cy = diag_top + diag_h/2 - bh/2

    node_data = [
        (sx,                        cy, PALE, LGRAY, NAVY, "🏢", "物件オーナー", "物件提供\n管理委託"),
        (sx + (bw+Inches(0.6)),     cy, NAVY, BLUE,  WHITE,"⚡", "SHIKI×日本住戸\nPlatform", "月次課金管理\nサービス統合"),
        (sx + 2*(bw+Inches(0.6)),   cy, PALE, LGRAY, NAVY, "🏪", "仲 介 会 社", "入居者紹介\nストック収益"),
        (sx + 3*(bw+Inches(0.6)),   cy, PALE, LGRAY, NAVY, "👤", "入 居 者",   "月額課金\nサービス享受"),
    ]

    for (bx, by, bg, bc, tc, icon, title, sub) in node_data:
        add_rect(slide, bx, by, bw, bh, fill=bg, line=bc, radius=30000)
        sc = RGBColor(0xAA,0xCC,0xEE) if bg==NAVY else GRAY
        add_textbox(slide, bx, by+Inches(0.1), bw, Inches(0.38),
                    icon, size=22, align=PP_ALIGN.CENTER)
        add_textbox(slide, bx, by+Inches(0.46), bw, Inches(0.44),
                    title, size=10, bold=True, color=tc, align=PP_ALIGN.CENTER)
        add_textbox(slide, bx, by+Inches(0.9), bw, Inches(0.56),
                    sub, size=9, color=sc, align=PP_ALIGN.CENTER)

    # Arrows between boxes
    arrow_labels = [
        ("委託", "収益還元"),
        ("紹介", "手数料"),
        ("マッチング", "月500円〜"),
    ]
    for i in range(3):
        bx_l = node_data[i][0]
        bx_r = node_data[i+1][0]
        ax   = bx_l + bw + Inches(0.06)
        ay   = cy + bh/2
        gap_w = bx_r - bx_l - bw - Inches(0.06)
        conn = slide.shapes.add_connector(1, ax, ay, ax+gap_w-Inches(0.06), ay)
        conn.line.color.rgb = BLUE
        conn.line.width = Pt(1.5)
        lbl, sub = arrow_labels[i]
        add_textbox(slide, ax, ay-Inches(0.28), gap_w-Inches(0.06), Inches(0.22),
                    lbl, size=8, color=GRAY, align=PP_ALIGN.CENTER)
        add_textbox(slide, ax, ay+Inches(0.06), gap_w-Inches(0.06), Inches(0.22),
                    sub, size=8, color=BLUE, align=PP_ALIGN.CENTER)

    # Attached services row at bottom
    svc_y = cy + bh + Inches(0.3)
    svc_items = [("⚡ 電気", RGBColor(0xFF,0xF3,0xE0), ACCENT),
                 ("🎬 コンテンツ", RGBColor(0xF3,0xE5,0xF5), RGBColor(0x90,0x60,0xB8)),
                 ("📡 ネット", RGBColor(0xE3,0xF2,0xFD), RGBColor(0x20,0x80,0xC0))]
    sw2 = Inches(1.3)
    sx2 = node_data[1][0] + (bw - (len(svc_items)*sw2 + Inches(0.1)*2)) / 2
    for k, (lbl, bg2, fg2) in enumerate(svc_items):
        xi = sx2 + k*(sw2+Inches(0.1))
        add_rect(slide, xi, svc_y, sw2, Inches(0.34), fill=bg2, line=fg2, radius=20000)
        add_textbox(slide, xi, svc_y+Inches(0.06), sw2, Inches(0.26),
                    lbl, size=9, color=fg2, align=PP_ALIGN.CENTER)
        # dashed connector
        conn2 = slide.shapes.add_connector(1, xi+sw2/2, cy+bh, xi+sw2/2, svc_y)
        conn2.line.color.rgb = GRAY
        conn2.line.width = Pt(0.75)
        conn2.line.dash_style = 4   # dash

    # Curved return arrow (入居者 → Platform) label
    add_textbox(slide, node_data[1][0]+bw/2-Inches(1), diag_top+Inches(0.15),
                Inches(2), Inches(0.26),
                "← 月額料金支払い", size=8, color=ACCENT, align=PP_ALIGN.CENTER)

    return slide

# ═══════════════════════════════════════════════════════
# SLIDE 7 – 手数料
# ═══════════════════════════════════════════════════════
def slide_fee():
    slide = prs.slides.add_slide(BLANK)
    add_header(slide, "07 / Fee Structure", "手 数 料 イ メ ー ジ", "ストック収益モデル")

    body_top = Inches(0.97)
    body_h   = H - body_top - Inches(0.2)
    cw = (W - Inches(0.35*2 + 0.28)) / 2
    L_x = Inches(0.35)
    R_x = L_x + cw + Inches(0.28)

    # ── 左: 手数料内訳 ──────────────────────────────────
    add_card(slide, L_x, body_top, cw, body_h, fill=PALE, border=LGRAY)
    tx = L_x + Inches(0.22)
    ty = body_top + Inches(0.18)
    lw = cw - Inches(0.44)

    add_textbox(slide, tx, ty, lw, Inches(0.26),
                "手数料内訳イメージ", size=9, bold=True, color=GRAY)
    ty += Inches(0.36)

    fee_rows = [
        ("入居者月額基本料", "¥500 / 月"),
        ("電気サービス収益", "●% シェア"),
        ("コンテンツ収益",  "●% シェア"),
        ("インターネット収益","●% シェア"),
    ]
    for name, val in fee_rows:
        add_rect(slide, tx, ty+Inches(0.26), lw, Pt(0.75), fill=LGRAY)
        add_textbox(slide, tx, ty, lw*0.55, Inches(0.28), name, size=10, color=NAVY)
        add_textbox(slide, tx+lw*0.55, ty, lw*0.45, Inches(0.28),
                    val, size=12, bold=True, color=BLUE, align=PP_ALIGN.RIGHT)
        ty += Inches(0.42)

    # total row
    add_rect(slide, tx, ty, lw, Pt(2), fill=BLUE)
    ty += Inches(0.12)
    add_textbox(slide, tx, ty, lw*0.6, Inches(0.38),
                "仲介会社への還元率", size=11, bold=True, color=NAVY)
    add_textbox(slide, tx+lw*0.6, ty, lw*0.4, Inches(0.38),
                "●%", size=22, bold=True, color=BLUE, align=PP_ALIGN.RIGHT)
    ty += Inches(0.55)

    # note
    note = add_rect(slide, tx, ty, lw, Inches(0.66), fill=WHITE, line=LGRAY, radius=15000)
    add_textbox(slide, tx+Inches(0.12), ty+Inches(0.1), lw-Inches(0.24), Inches(0.5),
                "※ 入居者が在籍する限り毎月継続して受け取れるストック型収益。\n　 戸数が増えるほど積み上がります。",
                size=8, color=GRAY)

    # ── 右: 棒グラフ ────────────────────────────────────
    add_card(slide, R_x, body_top, cw, body_h, fill=PALE, border=LGRAY)
    tx2 = R_x + Inches(0.22)
    ty2 = body_top + Inches(0.18)
    lw2 = cw - Inches(0.44)

    add_textbox(slide, tx2, ty2, lw2, Inches(0.26),
                "ストック収益積み上がりイメージ", size=9, bold=True, color=GRAY)
    ty2 += Inches(0.36)

    # bar chart
    chart_data = [1, 2, 3.2, 4.5, 6.2, 8, 9.8, 11.5, 13.5, 15.2, 17, 19]
    chart_h = Inches(2.5)
    chart_w = lw2
    chart_y = ty2 + Inches(0.1)
    max_v   = max(chart_data)
    n_bars  = len(chart_data)
    bar_w   = (chart_w - Inches(0.1)*(n_bars-1)) / n_bars

    for j, v in enumerate(chart_data):
        bh_px = chart_h * (v / max_v)
        bx_j  = tx2 + j*(bar_w + Inches(0.1))
        by_j  = chart_y + chart_h - bh_px

        # gradient simulation: darker at bottom
        add_rect(slide, bx_j, by_j, bar_w, bh_px, fill=BLUE, radius=8000)

    # x-axis
    add_rect(slide, tx2, chart_y+chart_h+Inches(0.04), lw2, Pt(1), fill=LGRAY)

    x_labels = ["1", "2", "3", "4", "5", "6", "9", "12", "15", "18", "21", "24ヶ月"]
    for j, lbl in enumerate(x_labels):
        bx_j = tx2 + j*(bar_w + Inches(0.1))
        add_textbox(slide, bx_j, chart_y+chart_h+Inches(0.08),
                    bar_w, Inches(0.2), lbl, size=6, color=GRAY, align=PP_ALIGN.CENTER)

    ty2 = chart_y + chart_h + Inches(0.38)
    add_textbox(slide, tx2, ty2, lw2, Inches(0.2),
                "※ 保有戸数・稼働率に応じた試算イメージ。実際の収益は条件によって異なります。",
                size=7, color=GRAY)
    ty2 += Inches(0.3)

    # 2 summary boxes
    sw3 = (lw2 - Inches(0.14)) / 2
    for k, (unit, val, col) in enumerate([
        ("100戸×24ヶ月 目安", "●●万円/月", BLUE),
        ("500戸×24ヶ月 目安", "●●万円/月", ACCENT),
    ]):
        sx3 = tx2 + k*(sw3+Inches(0.14))
        add_card(slide, sx3, ty2, sw3, Inches(0.72), fill=WHITE, border=LGRAY)
        add_textbox(slide, sx3, ty2+Inches(0.06), sw3, Inches(0.2),
                    unit, size=7, color=GRAY, align=PP_ALIGN.CENTER)
        add_textbox(slide, sx3, ty2+Inches(0.3), sw3, Inches(0.34),
                    val, size=14, bold=True, color=col, align=PP_ALIGN.CENTER)

    return slide

# ═══════════════════════════════════════════════════════
# SLIDE 8 – 問い合わせ先
# ═══════════════════════════════════════════════════════
def slide_contact():
    slide = prs.slides.add_slide(BLANK)
    add_header(slide, "08 / Contact", "問 い 合 わ せ 先")

    body_top = Inches(0.97)
    body_h   = H - body_top - Inches(0.2)
    cw = (W - Inches(0.35*2 + 0.28)) / 2
    L_x = Inches(0.35)
    R_x = L_x + cw + Inches(0.28)

    def contact_card(cx, brand, brand_color, rows):
        add_card(slide, cx, body_top, cw, body_h - Inches(0.66), fill=PALE, border=LGRAY)
        tx = cx + Inches(0.3)
        ty = body_top + Inches(0.22)
        lw = cw - Inches(0.6)
        add_textbox(slide, tx, ty, lw, Inches(0.42),
                    brand, size=18, bold=True, color=brand_color)
        ty += Inches(0.52)
        add_rect(slide, tx, ty, lw, Pt(1), fill=LGRAY)
        ty += Inches(0.16)
        for lbl, val in rows:
            add_textbox(slide, tx, ty, lw, Inches(0.2), lbl, size=7, color=GRAY)
            ty += Inches(0.2)
            add_textbox(slide, tx, ty, lw, Inches(0.3), val, size=11, bold=True, color=NAVY)
            ty += Inches(0.34)
        ty += Inches(0.1)
        add_textbox(slide, tx, ty, lw, Inches(0.22),
                    "受付時間: 平日 ●:00〜●:00", size=8, color=GRAY)

    contact_card(L_x, "SHIKI", BLUE, [
        ("担当部署", "●●●●部"),
        ("担当者名", "●●　●●"),
        ("電話番号", "0●-●●●●-●●●●"),
        ("メールアドレス", "●●●●@shiki.co.jp"),
    ])
    contact_card(R_x, "日本住戸", NAVY, [
        ("担当部署", "●●●●部"),
        ("担当者名", "●●　●●"),
        ("電話番号", "0●-●●●●-●●●●"),
        ("メールアドレス", "●●●●@nihon-juto.co.jp"),
    ])

    # footer banner
    footer_y = H - Inches(0.78)
    foot = add_rect(slide, Inches(0.35), footer_y, W-Inches(0.7), Inches(0.56),
                    fill=NAVY, radius=15000)
    add_textbox(slide, Inches(0.6), footer_y+Inches(0.12), W-Inches(1.2), Inches(0.34),
                "お気軽にお問い合わせください　|　SHIKI（親会社）× 日本住戸（グループ会社）",
                size=10, color=RGBColor(0xAA,0xBB,0xCC), align=PP_ALIGN.CENTER)

    return slide

# ═══════════════════════════════════════════════════════
# Build all slides
# ═══════════════════════════════════════════════════════
slide_company_overview()
slide_business_model()
slide_broker_network()
slide_product()
slide_services()
slide_proposal()
slide_fee()
slide_contact()

out = "/home/user/contact-form/presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
